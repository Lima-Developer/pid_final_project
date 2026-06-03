from __future__ import annotations

from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
PRESENTATION_DIR = ROOT / "presentation"
RESULTS_DIR = ROOT / "results"
DATA_DIR = ROOT / "data" / "single_case"
OUTPUT_PPTX = PRESENTATION_DIR / "Face_Privacy_Blur.pptx"

BG = RGBColor(245, 241, 234)
SURFACE = RGBColor(255, 255, 255)
INK = RGBColor(31, 36, 48)
MUTED = RGBColor(88, 96, 116)
LINE = RGBColor(220, 214, 203)
ACCENT = RGBColor(0, 145, 105)
WARNING = RGBColor(224, 93, 83)
DARK = RGBColor(22, 27, 34)


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left: float, top: float, width: float, height: float, text: str, size: int = 14, color: RGBColor = MUTED, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold


def add_kicker(slide, text: str) -> None:
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.58), Inches(0.12), Inches(0.15), Inches(0.15))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT
    marker.line.color.rgb = ACCENT
    add_text(slide, 0.82, 0.1, 6.0, 0.22, text.upper(), size=8, color=MUTED, bold=True)


def add_title(slide, title: str, subtitle: str = "") -> None:
    add_text(slide, 0.55, 0.28, 12.0, 0.55, title, size=25, color=INK, bold=True)
    if subtitle:
        add_text(slide, 0.57, 0.82, 12.1, 0.38, subtitle, size=11, color=MUTED)


def add_card(slide, left: float, top: float, width: float, height: float, title: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    if title:
        add_text(slide, left + 0.18, top + 0.1, width - 0.36, 0.28, title, size=12, color=INK, bold=True)
    return shape


def add_bullets(slide, left: float, top: float, width: float, height: float, title: str, bullets: list[str], size: int = 12) -> None:
    add_card(slide, left, top, width, height, title)
    box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.5), Inches(width - 0.44), Inches(height - 0.6))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = MUTED


def add_formula(slide, left: float, top: float, width: float, height: float, title: str, formula: str, explanation: str) -> None:
    add_card(slide, left, top, width, height, title)
    add_text(slide, left + 0.25, top + 0.58, width - 0.5, 0.48, formula, size=17, color=INK, bold=True)
    add_text(slide, left + 0.25, top + 1.12, width - 0.5, height - 1.25, explanation, size=11, color=MUTED)


def add_picture(slide, image_path: Path, left: float, top: float, width: float, height: float, title: str | None = None) -> None:
    add_card(slide, left, top, width, height, title)
    image = cv2.imread(str(image_path))
    if image is None:
        add_text(slide, left + 0.25, top + 0.65, width - 0.5, height - 0.8, f"Imagem não encontrada: {image_path.name}", size=12, color=WARNING)
        return

    image_height, image_width = image.shape[:2]
    box_left = left + 0.12
    box_top = top + (0.5 if title else 0.12)
    box_width = width - 0.24
    box_height = height - (0.64 if title else 0.24)
    image_ratio = image_width / image_height
    box_ratio = box_width / box_height

    if image_ratio > box_ratio:
        render_width = box_width
        render_height = box_width / image_ratio
    else:
        render_height = box_height
        render_width = box_height * image_ratio

    render_left = box_left + (box_width - render_width) / 2
    render_top = box_top + (box_height - render_height) / 2
    slide.shapes.add_picture(str(image_path), Inches(render_left), Inches(render_top), Inches(render_width), Inches(render_height))


def new_slide(prs: Presentation, kicker: str, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, kicker)
    add_title(slide, title, subtitle)
    return slide


def build_slides(prs: Presentation) -> None:
    source_image = DATA_DIR / "street_crowd.jpg"
    final_panel = RESULTS_DIR / "final_panel.png"
    final_result = RESULTS_DIR / "final_result.png"
    face_mask = RESULTS_DIR / "face_mask.png"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK)
    add_text(slide, 0.7, 0.55, 6.0, 0.35, "PROCESSAMENTO DIGITAL DE IMAGENS", size=10, color=RGBColor(180, 190, 205), bold=True)
    add_text(slide, 0.7, 1.1, 6.2, 1.0, "Face Privacy Blur", size=38, color=RGBColor(255, 255, 255), bold=True)
    add_text(slide, 0.72, 2.0, 5.7, 0.95, "Detecção de rostos e desfoque localizado para privacidade em uma única imagem urbana real.", size=15, color=RGBColor(215, 220, 230))
    add_picture(slide, final_result, 6.35, 0.72, 6.55, 5.75, "Resultado final")
    add_text(slide, 0.72, 6.55, 6.2, 0.25, "Fonte da imagem: Wikimedia Commons — Street Crowd.jpg — CC0", size=9, color=RGBColor(170, 180, 195))

    slide = new_slide(prs, "CASO ESCOLHIDO", "O foco principal é privacidade facial", "O projeto se enquadra diretamente no exemplo: detectar rostos e aplicar desfoque para privacidade.")
    add_bullets(slide, 0.65, 1.4, 4.0, 4.7, "Problema realista", [
        "Cenas urbanas podem expor rostos de pessoas.",
        "Privacidade exige reduzir informação identificável.",
        "O escopo é pequeno: uma imagem, um pipeline, um painel.",
        "As demais técnicas entram como análise e apoio.",
    ])
    add_picture(slide, source_image, 4.95, 1.4, 7.75, 4.7, "Imagem única usada no projeto")

    slide = new_slide(prs, "REQUISITOS", "Um caso principal e três complementos do enunciado", "A defesa deve deixar claro que o caso escolhido é privacidade facial.")
    add_bullets(slide, 0.65, 1.35, 3.0, 4.9, "Caso principal", [
        "Detectar rostos.",
        "Aplicar desfoque.",
        "Preservar privacidade.",
        "Não reconhecer identidade.",
    ])
    add_bullets(slide, 3.95, 1.35, 3.0, 4.9, "PDI clássica", [
        "Blur.",
        "Threshold.",
        "Sobel/Canny.",
        "Harris/ORB.",
    ])
    add_bullets(slide, 7.25, 1.35, 3.0, 4.9, "Composição", [
        "ROI facial.",
        "Máscara.",
        "Blend.",
        "Paste.",
    ])
    add_bullets(slide, 10.55, 1.35, 2.1, 4.9, "YOLO", [
        "Pessoa.",
        "Caixa.",
        "Classe.",
        "Score.",
    ], size=11)

    slide = new_slide(prs, "MODELO MATEMÁTICO", "Uma imagem digital é uma matriz de pixels", "A base do projeto é transformar e interpretar matrizes.")
    add_formula(slide, 0.65, 1.35, 3.95, 2.25, "Imagem RGB", "I(x,y) = [R,G,B]", "Cada pixel tem posição e três canais de cor. O processamento transforma esses valores para extrair informação.")
    add_formula(slide, 4.85, 1.35, 3.95, 2.25, "Escala de cinza", "Y = 0,299R + 0,587G + 0,114B", "A luminância reduz a imagem para uma matriz e aproxima a sensibilidade visual humana.")
    add_formula(slide, 9.05, 1.35, 3.05, 2.25, "ROI facial", "B = (x1,y1,x2,y2)", "Cada rosto detectado vira uma submatriz retangular da imagem.")
    add_picture(slide, source_image, 1.15, 4.05, 11.1, 2.25, "Imagem como dados numéricos")

    slide = new_slide(prs, "DETECÇÃO FACIAL", "Haar Cascade procura padrões de contraste parecidos com rosto", "O detector clássico localiza regiões candidatas antes do blur.")
    add_formula(slide, 0.65, 1.35, 4.0, 2.45, "Haar-like features", "f = soma(A) − soma(B)", "O detector compara somas de intensidades em regiões claras e escuras. Isso captura padrões como olhos, nariz e testa.")
    add_formula(slide, 4.95, 1.35, 4.0, 2.45, "Integral image", "S(x,y)=Σ I(i,j)", "A imagem integral acelera o cálculo de somas retangulares, permitindo varrer a imagem em várias escalas.")
    add_bullets(slide, 9.25, 1.35, 3.1, 2.45, "No projeto", [
        "8 rostos/possíveis rostos.",
        "Caixas azuis no quadro 1.",
        "Base para a máscara facial.",
    ], size=10)
    add_picture(slide, final_panel, 0.9, 4.1, 11.8, 2.25, "No painel: quadro 1 mostra rostos detectados")

    slide = new_slide(prs, "PRÉ-PROCESSAMENTO", "Blur reduz ruído; threshold separa regiões por intensidade", "Essa etapa ajuda a interpretar a imagem antes da privacidade final.")
    add_formula(slide, 0.65, 1.35, 4.0, 2.15, "Gaussian Blur", "I' = Gσ * I", "Convolução: cada pixel vira média ponderada da vizinhança, reduzindo ruídos locais.")
    add_formula(slide, 4.95, 1.35, 4.0, 2.15, "Median Blur", "I'(x,y) = mediana(vizinhança)", "Filtro robusto contra valores extremos, útil para suavizar sem destruir totalmente bordas.")
    add_formula(slide, 9.25, 1.35, 3.2, 2.15, "Otsu", "t* = argmin σ²w(t)", "Escolhe automaticamente o limiar que separa melhor claro e escuro.")
    add_picture(slide, final_panel, 0.9, 3.85, 11.8, 2.6, "No painel: quadro 2 mostra blur + threshold")

    slide = new_slide(prs, "BORDAS", "Bordas aparecem onde a intensidade muda rapidamente", "A matemática aqui é derivada discreta sobre a matriz de pixels.")
    add_formula(slide, 0.65, 1.35, 4.0, 2.35, "Sobel", "|∇I| = √(Ix² + Iy²)", "Sobel estima derivadas horizontal e vertical. Quanto maior o gradiente, mais provável é uma borda.")
    add_formula(slide, 4.95, 1.35, 4.0, 2.35, "Canny", "blur → gradiente → supressão → histerese", "Canny combina suavização, direção do gradiente, afinamento e limiares duplos.")
    add_bullets(slide, 9.25, 1.35, 3.1, 2.35, "Aplicação", [
        "Contornos de rostos e corpos.",
        "Separação roupa/fundo.",
        "Estrutura visual da multidão.",
    ], size=11)
    add_picture(slide, final_panel, 0.9, 4.05, 11.8, 2.35, "No painel: quadro 3 mostra as bordas Canny")

    slide = new_slide(prs, "CANTOS E PONTOS", "Harris e ORB mostram regiões visualmente informativas", "Essas técnicas não desfocam rostos; elas explicam a estrutura da cena.")
    add_formula(slide, 0.65, 1.35, 4.4, 2.5, "Harris", "R = det(M) − k·trace(M)²", "A matriz M resume variações locais. Quando há variação em duas direções, a região tende a ser canto.")
    add_formula(slide, 5.35, 1.35, 3.5, 2.5, "ORB", "keypoint + descriptor", "ORB detecta pontos robustos e descreve padrões locais com vetores binários.")
    add_bullets(slide, 9.15, 1.35, 3.2, 2.5, "No projeto", [
        "Texturas de roupa.",
        "Bordas de rosto.",
        "Letras e objetos.",
        "Pontos verdes no quadro 4.",
    ], size=10)
    add_picture(slide, final_panel, 0.9, 4.1, 11.8, 2.25, "No painel: quadro 4 mostra os pontos ORB")

    slide = new_slide(prs, "MÁSCARA E PRIVACIDADE", "O blur é aplicado só onde a máscara facial é positiva", "Essa é a etapa central do caso escolhido pelo grupo.")
    add_formula(slide, 0.65, 1.35, 4.0, 2.45, "Máscara binária", "M(x,y) ∈ {0,1}", "A máscara marca pixels dentro das regiões faciais. Fora da máscara, a imagem original é preservada.")
    add_formula(slide, 4.95, 1.35, 4.0, 2.45, "Composição", "Iout = M·Blur(I)+(1−M)·I", "A saída mistura imagem desfocada nas faces e imagem original no restante.")
    add_picture(slide, face_mask, 9.35, 1.35, 2.8, 2.45, "Máscara facial")
    add_picture(slide, final_result, 0.9, 4.1, 11.8, 2.25, "Resultado: faces desfocadas e cena preservada")

    slide = new_slide(prs, "ROI, BLEND E PASTE", "A composição visual demonstra manipulação local da imagem", "O requisito de ROI/máscara/blend/paste é atendido sobre as regiões faciais.")
    add_formula(slide, 0.65, 1.35, 4.0, 2.35, "ROI", "face = I[y1:y2,x1:x2]", "Cada rosto é extraído como subimagem e processado localmente.")
    add_formula(slide, 4.95, 1.35, 4.0, 2.35, "Paste", "I[35:185,35:185] ← ROI blur", "Uma ROI facial desfocada é ampliada e colada no canto como detalhe visual.")
    add_bullets(slide, 9.25, 1.35, 3.1, 2.35, "Por que isso importa?", [
        "Mostra ROI.",
        "Mostra máscara.",
        "Mostra blend.",
        "Mostra paste.",
    ], size=11)
    add_picture(slide, final_result, 0.9, 4.05, 11.8, 2.35, "Imagem final com inset da ROI desfocada")

    slide = new_slide(prs, "YOLO", "YOLO complementa o projeto detectando pessoas na mesma cena", "Assim também atendemos o caso de identificar objetos em ambiente externo.")
    add_bullets(slide, 0.65, 1.35, 4.0, 4.8, "No pipeline", [
        "Modelo ONNX pré-treinado.",
        "OpenCV DNN.",
        "Filtro da classe `person`.",
        "Caixas e scores de confiança.",
    ])
    add_formula(slide, 4.95, 1.35, 3.8, 2.15, "Confiança", "score = P(objeto)·P(classe)", "A caixa é aceita se a confiança passar do limiar. Caixas redundantes são filtradas por NMS.")
    add_picture(slide, final_panel, 8.95, 1.35, 3.4, 4.8, "Quadro YOLO no painel")

    slide = new_slide(prs, "PAINEL FINAL", "Um único painel mostra todo o projeto", "Este é o principal artefato para a defesa.")
    add_picture(slide, final_panel, 0.55, 1.3, 12.2, 5.55, "Seis quadros: rostos, threshold, Canny, ORB, YOLO e privacidade")

    slide = new_slide(prs, "LIMITAÇÕES", "O sistema é uma demonstração controlada, não uma solução de produção", "Limitações bem explicadas deixam a defesa mais forte.")
    add_bullets(slide, 0.75, 1.35, 5.7, 4.9, "Limitações técnicas", [
        "Haar Cascade pode gerar falsos positivos.",
        "Rostos laterais ou ocluídos podem não ser detectados.",
        "YOLO detecta pessoas, não necessariamente faces.",
        "Iluminação e escala afetam a detecção.",
    ])
    add_bullets(slide, 6.75, 1.35, 5.7, 4.9, "Limitações éticas", [
        "Não é reconhecimento facial.",
        "Não identifica nome ou documento.",
        "Reduz exposição visual, não garante anonimato perfeito.",
        "Uso real exigiria validação mais rigorosa.",
    ])

    slide = new_slide(prs, "DEMONSTRAÇÃO", "Plano de fala para defender com proficiência", "Use esta sequência na apresentação ao vivo.")
    add_bullets(slide, 0.75, 1.35, 3.7, 4.9, "1. Rodar", [
        "`python -m src.main`",
        "Mostrar que gera painel único.",
        "Se YOLO falhar: `--skip-yolo`.",
    ], size=11)
    add_bullets(slide, 4.65, 1.35, 3.7, 4.9, "2. Explicar", [
        "Imagem como matriz.",
        "Face ROI.",
        "Máscara binária.",
        "Blur localizado.",
        "Bordas/cantos e YOLO.",
    ], size=11)
    add_bullets(slide, 8.55, 1.35, 3.7, 4.9, "3. Fechar", [
        "Caso principal atendido.",
        "Complementos atendidos.",
        "Matemática aplicada.",
        "Limitações honestas.",
    ], size=11)

    slide = new_slide(prs, "CONCLUSÃO", "O projeto atende o caso de privacidade facial e ainda cobre os complementos", "A força da entrega é conectar matemática, código e resultado visual.")
    add_bullets(slide, 0.8, 1.4, 5.6, 4.7, "Mensagem final", [
        "O objetivo principal é detectar rostos e desfocá-los.",
        "As técnicas clássicas explicam a estrutura visual.",
        "YOLO detecta pessoas como objeto complementar.",
        "ROI, máscara, blend e paste aparecem no resultado final.",
    ])
    add_picture(slide, final_result, 6.85, 1.4, 5.65, 4.7, "Artefato final")
    add_text(slide, 0.85, 6.55, 11.8, 0.25, "Pergunta esperada: “isso reconhece pessoas?” Resposta: não; o objetivo é desfocar faces para privacidade.", size=10, color=WARNING, bold=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slides(prs)
    prs.save(OUTPUT_PPTX)
    print(f"Apresentação salva em: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
